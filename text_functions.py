from pptx import Presentation
from tkinter import messagebox
import openpyxl as xl
import os

path = os.getcwd()

def t_runner(xl_path, ppt_path):
    prs = Presentation(ppt_path)
    slides = [slide for slide in prs.slides]
    shapes = []
    for slide in slides:
        for shape in slide.shapes:
            shapes.append(shape)
    
    replace_dic = xl_to_dic(xl_path, "Sheet1")
    replace_text(replace_dic, shapes)
    try:
        prs.save(os.path.join(path, "output.pptx"))
        messagebox.showinfo("Info", "Operation Succeeded")
    except PermissionError:
        messagebox.showerror(title="Error", message="Please Close the output file and try again")
    except:
        messagebox.showerror(title="Erro", message="Something else went wrong")

def replace_text(replacements, shapes):
    """Takes dict of {match: replacement, ... } and replaces all matches.
    Currently not implemented for charts or graphics.
    """
    for shape in shapes:
        for match, replacement in replacements.items():
            if shape.has_text_frame:
                if (shape.text.find(match)) != -1:
                    text_frame = shape.text_frame
                    for paragraph in text_frame.paragraphs:
                        whole_text = "".join(run.text for run in paragraph.runs)
                        whole_text = whole_text.replace(str(match), str(replacement))
                        for idx, run in enumerate(paragraph.runs):
                            if idx != 0:
                                p = paragraph._p
                                p.remove(run._r)
                        if(not(not paragraph.runs)):
                            paragraph.runs[0].text = whole_text


def xl_to_dic(path, sheet_name):

    wb = xl.load_workbook(path)

    sheet = wb[sheet_name]

    dic = {}

    for i in range(1, sheet.max_row + 1):
        if(sheet.cell(row=i, column=2).value == None):
            dic[(sheet.cell(row=i, column=1).value).rstrip()] = ""
        else:
            dic[(sheet.cell(row=i, column=1).value).rstrip()] = sheet.cell(row=i, column=2).value

    return dic