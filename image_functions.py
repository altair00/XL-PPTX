from pptx import Presentation
from tkinter import messagebox
import openpyxl as xl
from openpyxl_image_loader import SheetImageLoader
from openpyxl.utils import get_column_letter
from os import path
import os
import random
import shutil

main_path = os.getcwd()
image_path = os.path.join(main_path, '.image')

#image_top execution command coming directly to i_runner
#main method
def i_runner(xl_path, ppt_path):
#checking if the .image path exists or not
#if not exists try to make one and catch the exceptions(eg: permission error)
    if(path.exists(image_path)):
        shutil.rmtree(image_path)
    try:
        os.mkdir(image_path)
    except PermissionError:
        messagebox.showerror(title="Error", message="Please Try Again")
    wb = xl.load_workbook(xl_path)
    sheet = wb['Sheet1']
#sending a random number for saving purpose
    image_dic = image_return(random.randint(100, 1000), sheet)
    prs = Presentation(ppt_path)
    slides = [slide for slide in prs.slides]


    for slide in slides:
        shapes = []
        for shape in slide.shapes:
            shapes.append(shape)
        for match, replacement in image_dic.items(): 
            shape_index = find_image_name(shapes, match)
            if(shape_index != -1):
                img = slide.shapes[shape_index]
                inmage = os.path.join(image_path, replacement)
                replace_img_slide(slide, img, inmage)

    try:
        prs.save(os.path.join(main_path, "output.pptx"))
        messagebox.showinfo("Info", "Operation Succeeded")
    except PermissionError:
        messagebox.showerror(title="Error", message="Please Close the output file and try again")
    except:
        messagebox.showerror(title="Erro", message="Something else went wrong")

#Steering method it will return a dictionary

def image_return(u_id, sheet):
    image_loader = SheetImageLoader(sheet)

    dic = {}

    for i in range(1, sheet.max_row+1):
        if(image_loader.image_in(get_column_letter(2)+str(i))):
            image = image_loader.get(get_column_letter(2)+str(i))
            name = str(u_id)+"_"+str(i)+".jpg"
            dic[(sheet.cell(row=i, column=1).value).rstrip()] = name
            image_save(image, name)


    return dic

#just save the image if not RGB convert it to RGB

def image_save(image, name):
    if(image.mode in ("RGBA", "P")):
        image = image.convert("RGB")
    image.save(os.path.join(image_path, name))

#using the dictionary and finding the matching name image in the shapes object
#else return -1 
def find_image_name(shapes, img_name):
    for index, shape in enumerate(shapes):
        if shape.shape_type == 13:
            if(shape._pic.nvPicPr.cNvPr.get('name') == img_name):
                return index
    return -1

#replace the appropriate image in the slide
def replace_img_slide(slide, img, img_path):
    # Replace the picture in the shape object (img) with the image in img_path.
    imgPic = img._pic
    imgRID = imgPic.xpath('./p:blipFill/a:blip/@r:embed')[0]
    imgPart = slide.part.related_parts[imgRID]
    with open(img_path, 'rb') as f:
        rImgBlob = f.read()
    # replace
    imgPart._blob = rImgBlob
